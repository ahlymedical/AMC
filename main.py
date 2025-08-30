import os
import io
import traceback
import time
import google.generativeai as genai
import docx 
from docx.document import Document as DocxDocument
from docx.text.paragraph import Paragraph
from docx.table import _Cell
import PyPDF2
from PIL import Image
from pptx import Presentation
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename

# --- إعداد التطبيق ---
app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app)

# --- تهيئة نموذج Gemini ---
model = None
api_key_error = None
try:
    GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
    if not GEMINI_API_KEY:
        raise ValueError("خطأ فادح: متغير البيئة GEMINI_API_KEY غير موجود أو فارغ. يرجى إضافته إلى إعدادات الخادم.")
    
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-1.5-pro-latest') 
    print("✅ تم إعداد Gemini API بنجاح باستخدام نموذج gemini-1.5-pro-latest.")

except Exception as e:
    api_key_error = str(e)
    print(f"!!!!!! خطأ فادح أثناء إعداد Gemini API: {api_key_error}")

# --- دوال الترجمة الأساسية ---

def translate_text_api(text_to_translate, target_lang):
    if not text_to_translate or not text_to_translate.strip():
        return ""
    try:
        prompt = f"""As a master translator, provide a professional and context-aware translation of the following text into {target_lang}. 
Your output must be ONLY the translated text, preserving the original tone and meaning. 
Do not add any extra explanations or introductory phrases.

Original Text:
\"\"\"{text_to_translate}\"\"\""""
        
        print("      - [API Call] إرسال النص إلى Gemini API...")
        start_time = time.time()
        response = model.generate_content(prompt)
        end_time = time.time()
        print(f"      - [API Success] تم استلام الرد من Gemini API في {end_time - start_time:.2f} ثانية.")
        
        return response.text.strip()
    except Exception as e:
        print(f"      - !!!!! [API Error] حدث خطأ أثناء الاتصال بـ Gemini API: {e}")
        return text_to_translate

def translate_docx_in_place(doc: DocxDocument, target_lang: str):
    print("  [DOCX] بدء الترجمة المباشرة لملف DOCX...")
    
    for i, para in enumerate(doc.paragraphs):
        for j, run in enumerate(para.runs):
            if run.text.strip():
                print(f"    - ترجمة الفقرة {i+1}، الجزء {j+1}...")
                run.text = translate_text_api(run.text, target_lang)

    for i, table in enumerate(doc.tables):
        for j, row in enumerate(table.rows):
            for k, cell in enumerate(row.cells):
                for l, para in enumerate(cell.paragraphs):
                    for m, run in enumerate(para.runs):
                        if run.text.strip():
                            print(f"    - ترجمة الجدول {i+1}, الصف {j+1}, الخلية {k+1}, الجزء {m+1}...")
                            run.text = translate_text_api(run.text, target_lang)
                        
    print("  [DOCX] انتهت الترجمة المباشرة لملف DOCX.")
    return doc

def translate_pptx_in_place(prs: Presentation, target_lang: str):
    print("  [PPTX] بدء الترجمة المباشرة لملف PowerPoint...")
    
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                for k, paragraph in enumerate(shape.text_frame.paragraphs):
                    for l, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            print(f"    - ترجمة الشريحة {i+1}, الشكل {j+1}, الجزء {l+1}...")
                            run.text = translate_text_api(run.text, target_lang)
            
            if shape.has_table:
                for k, row in enumerate(shape.table.rows):
                    for l, cell in enumerate(row.cells):
                        for m, paragraph in enumerate(cell.text_frame.paragraphs):
                            for n, run in enumerate(paragraph.runs):
                                if run.text.strip():
                                    print(f"    - ترجمة جدول الشريحة {i+1}, الصف {k+1}, الخلية {l+1}, الجزء {n+1}...")
                                    run.text = translate_text_api(run.text, target_lang)

    print("  [PPTX] انتهت الترجمة المباشرة لملف PowerPoint.")
    return prs

# --- دوال قراءة الملفات ---

def read_text_from_pdf(stream):
    print("  [PDF] بدء استخلاص النص من ملف PDF...")
    reader = PyPDF2.PdfReader(stream)
    text = '\n'.join([page.extract_text() for page in reader.pages if page.extract_text()])
    print("  [PDF] تم استخلاص النص من ملف PDF بنجاح.")
    return text

def create_docx_from_text(text):
    print("  [DOCX Create] بدء إنشاء ملف DOCX جديد من النص المترجم...")
    mem_file = io.BytesIO()
    doc = docx.Document()
    doc.add_paragraph(text)
    doc.save(mem_file)
    mem_file.seek(0)
    print("  [DOCX Create] تم إنشاء الملف بنجاح.")
    return mem_file

# --- مسارات التطبيق (API Endpoints) ---

@app.route('/')
def serve_index():
    return app.send_static_file('index.html')

@app.route('/translate-file', methods=['POST'])
def translate_file_handler():
    print("\n[Request Start] تم استلام طلب ترجمة ملف جديد.")
    start_time = time.time()

    if not model:
        print("!!!!!! [Request Error] فشل الطلب لأن نموذج Gemini غير مهيأ.")
        return jsonify({"error": f"خدمة الـ API غير مهيأة بشكل صحيح. السبب: {api_key_error}"}), 500
    if 'file' not in request.files:
        print("!!!!!! [Request Error] لم يتم العثور على ملف في الطلب.")
        return jsonify({"error": "لم يتم العثور على ملف في الطلب."}), 400
    file = request.files['file']
    if file.filename == '':
        print("!!!!!! [Request Error] لم يتم تحديد أي ملف.")
        return jsonify({"error": "لم يتم تحديد أي ملف."}), 400
    
    filename = secure_filename(file.filename)
    target_lang = request.form.get('target_lang', 'English')
    print(f"[File Info] اسم الملف: {filename}, اللغة المستهدفة: {target_lang}")
    
    try:
        mem_file = io.BytesIO()
        new_filename = f"translated_{os.path.splitext(filename)[0]}"
        mimetype = ''

        if filename.lower().endswith('.docx'):
            original_doc = docx.Document(file.stream)
            translated_doc = translate_docx_in_place(original_doc, target_lang)
            translated_doc.save(mem_file)
            new_filename += ".docx"
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'

        elif filename.lower().endswith('.pptx'):
            original_prs = Presentation(file.stream)
            translated_prs = translate_pptx_in_place(original_prs, target_lang)
            translated_prs.save(mem_file)
            new_filename += ".pptx"
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

        else:
            text_to_translate = ""
            if filename.lower().endswith('.pdf'):
                text_to_translate = read_text_from_pdf(file.stream)
            elif filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                print("  [Image] بدء استخلاص النص من الصورة...")
                image = Image.open(file.stream)
                prompt_parts = [f"Extract and translate the text in this image to {target_lang}. Provide only the translation.", image]
                response = model.generate_content(prompt_parts)
                text_to_translate = response.text
                print("  [Image] تم استخلاص النص من الصورة.")
            else:
                return jsonify({"error": "نوع الملف غير مدعوم."}), 400

            if not text_to_translate.strip():
                return jsonify({"error": "لم يتمكن النظام من استخلاص أي نص من الملف."}), 400
            
            print("  [Text] بدء ترجمة النص المستخلص...")
            translated_text = translate_text_api(text_to_translate, target_lang)
            print("  [Text] انتهت ترجمة النص المستخلص.")
            mem_file = create_docx_from_text(translated_text)
            new_filename += ".docx"
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'

        mem_file.seek(0)
        end_time = time.time()
        print(f"[Request Success] اكتملت معالجة الملف بنجاح في {end_time - start_time:.2f} ثانية. جاري إرسال الملف...")
        return send_file(mem_file, as_attachment=True, download_name=new_filename, mimetype=mimetype)

    except Exception as e:
        end_time = time.time()
        print(f"!!!!!! [Request Critical Error] حدث خطأ فادح بعد {end_time - start_time:.2f} ثانية: {e}")
        traceback.print_exc()
        return jsonify({"error": "حدث خطأ داخلي في الخادم أثناء معالجة الملف."}), 500

@app.route('/translate-text', methods=['POST'])
def translate_text_handler():
    if not model: return jsonify({"error": f"خدمة الـ API غير مهيأة بشكل صحيح. السبب: {api_key_error}"}), 500
    data = request.get_json()
    text = data.get('text')
    if not text: return jsonify({"error": "لم يتم توفير أي نص للترجمة."}), 400
    target_lang = data.get('target_lang', 'English')
    try:
        translated_text = translate_text_api(text, target_lang)
        return jsonify({"translated_text": translated_text})
    except Exception as e:
        print(f"!!!!!! حدث خطأ أثناء ترجمة النص: {e}")
        return jsonify({"error": "حدث خطأ داخلي أثناء عملية الترجمة."}), 500

# --- تشغيل التطبيق ---
if __name__ == "__main__":
    port = int(os.environ.get('PORT', 8080))
    app.run(debug=False, host='0.0.0.0', port=port)
